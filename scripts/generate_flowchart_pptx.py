"""Генерация блок-схемы Агент-Почта в PowerPoint (.pptx).

Запуск: python scripts/generate_flowchart_pptx.py
Выход: docs/AGENT-FLOWCHART.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "AGENT-FLOWCHART.pptx"

COLORS = {
    "start": (RGBColor(212, 237, 218), RGBColor(40, 167, 69)),
    "agent": (RGBColor(204, 229, 255), RGBColor(0, 123, 255)),
    "ud": (RGBColor(248, 215, 218), RGBColor(220, 53, 69)),
    "decision": (RGBColor(255, 243, 205), RGBColor(255, 193, 7)),
    "infra": (RGBColor(233, 236, 239), RGBColor(108, 117, 125)),
    "end_ok": (RGBColor(212, 237, 218), RGBColor(40, 167, 69)),
    "end_spam": (RGBColor(245, 245, 245), RGBColor(153, 153, 153)),
    "end_err": (RGBColor(255, 224, 224), RGBColor(220, 53, 69)),
}


def _set_shape_style(shape, fill_rgb: RGBColor, line_rgb: RGBColor, *, bold: bool = False) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(9)
        para.font.name = "Calibri"
        if bold:
            para.font.bold = True


def _add_box(slide, left, top, width, height, text: str, kind: str, *, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    fill, line = COLORS[kind]
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.text = text
    _set_shape_style(sh, fill, line, bold=kind.startswith("end") or kind == "start")
    return sh


def _add_diamond(slide, left, top, size, text: str):
    fill, line = COLORS["decision"]
    sh = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_DECISION, left, top, size, size * 0.72)
    sh.text = text
    _set_shape_style(sh, fill, line)
    return sh


def _arrow(slide, x1, y1, x2, y2) -> None:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = RGBColor(80, 80, 80)
    conn.line.width = Pt(1)


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.text = "Агент-Почта — блок-схема обработки входящей корреспонденции"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Calibri"

    legend = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.8))
    lt = legend.text_frame
    lt.word_wrap = True
    lines = [
        "Условные обозначения:",
        "",
        "Зелёный — старт / успешное завершение",
        "Синий — ИИ-агент (LangGraph + Celery)",
        "Жёлтый — решение (да / нет)",
        "Красный — сотрудник УД (human-in-the-loop)",
        "Серый — инфраструктура и внешние сервисы",
        "",
        "Пороги: SPAM_THRESHOLD=0.85 · GRAY=0.70 · DEPT=0.70",
        "UI: /agents/incoming-mail · API: :8080",
    ]
    for i, line in enumerate(lines):
        para = lt.paragraphs[0] if i == 0 else lt.add_paragraph()
        para.text = line
        para.font.size = Pt(14 if i == 0 else 12)
        para.font.name = "Calibri"
        if i == 0:
            para.font.bold = True


def _slide_infra_main(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = Inches(2.8), Inches(0.55)
    X = Inches(3.1)
    y = Inches(0.25)

    hdr = slide.shapes.add_textbox(Inches(0.3), y, Inches(9.4), Inches(0.35))
    hdr.text_frame.text = "Слайд 2 — Инфраструктура и основной поток (0–8)"
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.size = Pt(14)
    y += Inches(0.45)

    boxes = [
        ("infra", "0. Celery Beat (10 с)\nagent_pochta.poll_imap"),
        ("infra", "0.1 IMAP UNSEEN\nmail.turbo-don.ru:993"),
        ("infra", "0.2 process_email_task\nLangGraph"),
        ("start", "1. Старт — новое письмо"),
        ("agent", "2. imap_listener"),
        ("agent", "3. spam_filter — правила, Прил. А"),
    ]
    prev = None
    for kind, txt in boxes:
        sh = _add_box(slide, X, y, W, H, txt, kind)
        if prev:
            _arrow(slide, X + W / 2, prev.top + prev.height, X + W / 2, y)
        prev = sh
        y += H + Inches(0.12)

    d1 = _add_diamond(slide, X + Inches(0.55), y, Inches(1.7), "4. Спам\nпо правилам?")
    _arrow(slide, X + W / 2, prev.top + prev.height, X + W / 2, d1.top)
    y += Inches(1.35)

    spam_end = _add_box(
        slide, Inches(0.35), d1.top + Inches(0.05), Inches(2.2), Inches(0.55),
        "5A. Конец — spam", "end_spam",
    )
    main = _add_box(
        slide, X, y + Inches(0.2), W, Inches(0.85),
        "6–8. identify_sender (RAG)\nprocess_content (вложения)\nroute_department (1× LLM)", "agent",
    )
    _arrow(slide, d1.left + d1.width / 2, d1.top + d1.height, spam_end.left + spam_end.width, spam_end.top + spam_end.height / 2)
    _arrow(slide, X + W / 2, d1.top + d1.height, X + W / 2, main.top)

    note = slide.shapes.add_textbox(Inches(6.2), Inches(0.8), Inches(3.3), Inches(2.2))
    nt = note.text_frame
    nt.word_wrap = True
    for i, line in enumerate([
        "Сервисы:",
        "Qdrant — RAG",
        "LLM /chat/completions",
        "LocalDocumentService",
        "Integration Service → 1С",
    ]):
        p = nt.paragraphs[0] if i == 0 else nt.add_paragraph()
        p.text = line
        p.font.size = Pt(10)


def _slide_decisions_hitl(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = Inches(2.6), Inches(0.5)
    X = Inches(0.4)
    y = Inches(0.25)

    hdr = slide.shapes.add_textbox(X, y, Inches(9), Inches(0.35))
    hdr.text_frame.text = "Слайд 3 — Решения LLM и HITL (9–12)"
    hdr.text_frame.paragraphs[0].font.bold = True
    y += Inches(0.5)

    entry = _add_box(slide, X, y, W, H, "После LLM analyze_incoming", "agent")
    y += Inches(0.65)
    prev = entry
    for q, branch in [
        ("9. spam conf ≥ 0.85?", "Да → spam"),
        ("10. серая зона 0.70–0.85?", "Да → awaiting_human"),
        ("11. dept_conf < 0.70?", "Да → awaiting_human"),
    ]:
        d = _add_diamond(slide, X + Inches(0.35), y, Inches(1.9), q)
        _arrow(slide, prev.left + prev.width / 2, prev.top + prev.height, d.left + d.width / 2, d.top)
        side = slide.shapes.add_textbox(X + W + Inches(0.15), y + Inches(0.05), Inches(2.2), Inches(0.4))
        side.text_frame.text = branch
        side.text_frame.paragraphs[0].font.size = Pt(9)
        y += Inches(1.05)
        prev = d

    ok = _add_box(slide, X, y, W, H, "OK → create_erp_task", "agent")
    _arrow(slide, prev.left + prev.width / 2, prev.top + prev.height, ok.left + ok.width / 2, ok.top)

    hx, hy = Inches(5.0), Inches(0.9)
    _add_box(slide, hx, hy, Inches(4.5), Inches(0.4), "Петля 12.x — Сотрудник УД", "ud")
    hy += Inches(0.48)
    for s in [
        "GET …/email-messages?status=awaiting_human",
        "POST …/resolve-human (mark_spam / mark_not_spam / approve_routing)",
        "POST …/restore-from-spam",
        "Celery: continue_after_human · reprocess_message",
    ]:
        _add_box(slide, hx, hy, Inches(4.5), Inches(0.42), s, "ud")
        hy += Inches(0.48)


def _slide_erp_endpoints(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    X, y = Inches(0.4), Inches(0.25)
    W, H = Inches(3.0), Inches(0.55)

    hdr = slide.shapes.add_textbox(X, y, Inches(9), Inches(0.35))
    hdr.text_frame.text = "Слайд 4 — 1С, retry, API (13–15)"
    hdr.text_frame.paragraphs[0].font.bold = True
    y += Inches(0.5)

    prev = None
    for kind, txt, shape in [
        ("agent", "13. create_erp_task → 1С", MSO_SHAPE.ROUNDED_RECTANGLE),
        ("decision", "14. 1С OK?", MSO_SHAPE.FLOWCHART_DECISION),
        ("agent", "14.1 retry 600с × 5", MSO_SHAPE.ROUNDED_RECTANGLE),
        ("ud", "14.2 POST …/retry-erp", MSO_SHAPE.ROUNDED_RECTANGLE),
        ("end_ok", "15. done", MSO_SHAPE.ROUNDED_RECTANGLE),
        ("end_err", "15E. error → УД", MSO_SHAPE.ROUNDED_RECTANGLE),
    ]:
        sh = _add_box(slide, X, y, W, H, txt, kind, shape=shape)
        if prev:
            _arrow(slide, X + W / 2, prev.top + prev.height, X + W / 2, y)
        prev = sh
        y += H + Inches(0.14)

    api = slide.shapes.add_textbox(Inches(4.0), Inches(0.8), Inches(5.5), Inches(5.2))
    at = api.text_frame
    at.word_wrap = True
    for i, line in enumerate([
        "REST API (:8080):",
        "GET  /health",
        "GET  /api/v1/email-messages",
        "GET  /api/v1/email-messages/{id}",
        "POST /api/v1/email-messages/{id}/resolve-human",
        "POST /api/v1/email-messages/{id}/retry-erp",
        "POST /api/v1/email-messages/{id}/restore-from-spam",
        "",
        "Celery: poll_imap · process_email · continue_after_human · retry_erp",
        "Статусы: done · spam · awaiting_human · error",
    ]):
        p = at.paragraphs[0] if i == 0 else at.add_paragraph()
        p.text = line
        p.font.size = Pt(11 if i == 0 else 10)


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    _title_slide(prs)
    _slide_infra_main(prs)
    _slide_decisions_hitl(prs)
    _slide_erp_endpoints(prs)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
